#!/usr/bin/env python3
"""Deterministic QA linter for ccPaidMediaAnalysts deliverables.

Advisory only — this script never blocks and always exits 0. It encodes the
mechanical, zero-judgment checks documented in CLAUDE.md ("PPT Design
Conventions", "Output Format Conventions", "Session Learnings") so they don't
have to be re-implemented in a throwaway /tmp script on every run.

It is the deterministic layer beneath the `qa-reviewer` subagent: the subagent
runs this for the mechanical findings, then reasons about the things a regex
cannot decide (fabricated campaign labels, currency-footnote adequacy,
narrative-vs-numbers consistency).

Checks
------
PPTX decks (lint_deck):
  - forbidden electric-blue #003DA5 / alert-orange #FF6B35 anywhere
  - accent green #05DD4D used as a solid FILL (allowed only as a thin rule)
  - forbidden words "waste" / "consider" in any text frame or table cell
  - chart series missing data labels (number-format judged by the reviewer)
  - percent chart series stored as fractions but formatted with a literal
    "%" (e.g. 0.0"%") — renders ~100x too small; use a true percent operator
  - shapes out of slide bounds / heavy overlap between text-bearing shapes

Text outputs (.md / .csv) (lint_text_output):
  - forbidden hex / words
  - filename contains a YYYY-MM-DD date
  - data-source + agent attribution header present (reports)
  - currency-blend heuristic: CAD and USD in one file with no conversion note

Raw reconciliation helper (summarize_raw):
  - sums clicks / impressions / spend from a raw export, FILTERING the Bing
    "Total" footer row (the doubling bug in CLAUDE.md). Sniffs UTF-16LE vs
    utf-8-sig. Prints totals for the reviewer to compare against the report.

Keyword-dedup cross-check (lint_keyword_dedup, needs --raw):
  - flags any keyword-ADD recommendation (action_type contains ADD_KEYWORD)
    whose keyword is ALREADY an active keyword in the account — the
    duplicate/no-op trap. Builds the active set from the 'Search keyword'
    column of the raw export(s), UNIONED with any keyword a change-history
    export shows as added (pass the change-history CSV as an additional
    --raw) — this catches a keyword added before the performance export's
    window that still has zero impressions, which the Search-keyword-column
    check alone can never see. Exact normalized match only (close-variant
    judgement is left to the qa-reviewer). Never recommend re-adding a term
    the account already runs; carve to exact only when it is currently served
    via a broader match type. Always pass the change-history export via
    --raw alongside the performance export(s) when one is available.

Usage
-----
    python scripts/qa_lint.py outputs/reports/ppt/executive_summary_x_2026-05-24.pptx
    python scripts/qa_lint.py outputs/reports/google/search_audit_x_2026-05-24.md
    python scripts/qa_lint.py --raw data/raw/bing/his_search_terms.csv   # totals only
"""

from __future__ import annotations

import argparse
import csv
import io
import os
import re
import sys
from dataclasses import dataclass

# --- Rules sourced from CLAUDE.md (do not invent new ones here) ---
FORBIDDEN_HEX = {
    "003DA5": "electric blue — off-brand (use [Client] teal)",
    "FF6B35": "alert orange — not in the [Client] theme",
}
GREEN_FILL_HEX = "05DD4D"  # accent green — allowed only as a thin rule, never a solid fill
FORBIDDEN_WORD_PATTERNS = [
    (re.compile(r"\bwast\w*", re.I), "waste/wasted — reframe as efficiency/optimization"),
    (re.compile(r"\bconsider\b", re.I), "consider — use a directive verb instead"),
]
# Tolerance (EMU) before a shape is flagged as out of bounds (~0.03 in).
BOUNDS_TOLERANCE_EMU = 30000


@dataclass
class Finding:
    severity: str  # "block" | "nit" | "info"
    category: str
    location: str
    message: str

    def render(self) -> str:
        return f"  [{self.severity:5}] {self.category}: {self.location} — {self.message}"


# ----------------------------------------------------------------------------
# Deck linting
# ----------------------------------------------------------------------------
def _iter_shapes(shapes):
    """Yield every shape, descending into groups."""
    for shp in shapes:
        yield shp
        if shp.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            yield from _iter_shapes(shp.shapes)


def _shape_text(shape) -> str:
    parts = []
    if shape.has_text_frame:
        parts.append(shape.text_frame.text)
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                parts.append(cell.text)
    return "\n".join(p for p in parts if p)


def lint_deck(path: str) -> list[Finding]:
    findings: list[Finding] = []
    try:
        from pptx import Presentation
    except ImportError:
        return [Finding("info", "deps", path,
                        "python-pptx not installed — skipped deck lint "
                        "(pip install python-pptx)")]

    prs = Presentation(path)
    sw, sh = prs.slide_width, prs.slide_height

    for i, slide in enumerate(prs.slides, start=1):
        loc = f"slide {i}"
        # Forbidden hex / green-fill: scan the slide part XML AND each chart part
        # XML. PowerPoint stores chart series formatting in separate
        # ppt/charts/chart*.xml parts referenced from the slide, so colors there
        # never appear in slide._element.xml — scanning only the slide would let
        # an off-brand chart series slip through. Concatenating dedupes naturally
        # (one finding per hexcode per slide).
        xml = slide._element.xml
        for shp in _iter_shapes(slide.shapes):
            if getattr(shp, "has_chart", False):
                try:
                    xml += shp.chart._chartSpace.xml
                except (AttributeError, ValueError):
                    pass
        for hexcode, why in FORBIDDEN_HEX.items():
            if re.search(hexcode, xml, re.I):
                findings.append(Finding("block", "forbidden-color", loc,
                                        f"#{hexcode} present — {why}"))
        for m in re.finditer(r'<a:solidFill>\s*<a:srgbClr val="([0-9A-Fa-f]{6})"',
                             xml):
            if m.group(1).upper() == GREEN_FILL_HEX:
                findings.append(Finding("block", "forbidden-fill", loc,
                                        f"#{GREEN_FILL_HEX} used as a solid fill — "
                                        "accent green is allowed only as a thin rule"))
                break

        text_shapes = []
        for shp in _iter_shapes(slide.shapes):
            txt = _shape_text(shp)
            if txt.strip():
                text_shapes.append((shp, txt))
                for pat, why in FORBIDDEN_WORD_PATTERNS:
                    if pat.search(txt):
                        snippet = pat.search(txt).group(0)
                        findings.append(Finding("block", "forbidden-word", loc,
                                                f'"{snippet}" — {why}'))
            # Chart data labels
            if getattr(shp, "has_chart", False):
                if not _chart_has_labels(shp.chart):
                    findings.append(Finding("nit", "chart-labels", loc,
                                            "chart series missing data labels — add "
                                            "show_value with the right number_format"))
                findings.extend(_chart_pct_fraction_findings(shp.chart, loc))
            # Out-of-bounds geometry
            try:
                l, t, w, h = shp.left, shp.top, shp.width, shp.height
            except (AttributeError, TypeError):
                l = None
            if l is not None and None not in (t, w, h):
                if (l < -BOUNDS_TOLERANCE_EMU or t < -BOUNDS_TOLERANCE_EMU
                        or l + w > sw + BOUNDS_TOLERANCE_EMU
                        or t + h > sh + BOUNDS_TOLERANCE_EMU):
                    findings.append(Finding("nit", "geometry", loc,
                                            f"shape extends beyond slide bounds "
                                            f"(name={shp.name!r})"))

        # Heavy overlap between text-bearing shapes (signal-only: >60% of the
        # smaller shape's area covered).
        for a in range(len(text_shapes)):
            for b in range(a + 1, len(text_shapes)):
                ov = _overlap_ratio(text_shapes[a][0], text_shapes[b][0])
                if ov > 0.60:
                    findings.append(Finding("nit", "overlap", loc,
                                            f"text shapes overlap ~{ov:.0%} "
                                            f"({text_shapes[a][0].name!r} / "
                                            f"{text_shapes[b][0].name!r})"))
    return findings


def _chart_has_labels(chart) -> bool:
    """True if data labels are enabled at the plot OR series level.

    CLAUDE.md decks set labels per-series (`series.data_labels.show_value`), so a
    plot-level check alone yields false positives.
    """
    try:
        for plot in chart.plots:
            if plot.has_data_labels:
                return True
            for series in plot.series:
                try:
                    if series.data_labels.show_value:
                        return True
                except (ValueError, AttributeError):
                    continue
    except (ValueError, AttributeError):
        return True  # can't tell — don't cry wolf
    # XML fallback: chart-level <c:dLbls><c:showVal val="1"/> that the object
    # model doesn't surface as plot/series labels.
    try:
        if re.search(r'<c:showVal[^>]*val="1"', chart._chartSpace.xml):
            return True
    except (AttributeError, ValueError):
        pass
    return False


# Quoted percent inside a number-format string, e.g. 0.0"%"  — a LITERAL "%"
# character, NOT the percent operator. A literal "%" prints the stored number
# verbatim and appends a percent sign WITHOUT multiplying by 100, so a CTR
# stored as the fraction 0.0167 renders "0.0%" instead of "1.67%".
_LITERAL_PCT_RE = re.compile(r'"[^"]*%[^"]*"')
_CHART_NS = "http://schemas.openxmlformats.org/drawingml/2006/chart"


def _chart_pct_fraction_findings(chart, loc: str) -> list:
    """Flag a percent series stored as fractions but formatted with a literal "%".

    The trap: a series whose number format carries a quoted percent (0.0"%")
    while its values live in fraction space (all |v| < 1). PowerPoint then shows
    the number ~100x too small. The fix is a true percent operator (0.0%, which
    multiplies by 100) or percentage-point values (1.67) paired with the literal
    format. A genuinely sub-1% rate stored in percentage-point space (e.g. a
    [DSP Platform] display CTR of 0.07) also trips the |v| < 1 test, so the message
    hedges rather than asserting the values are definitely fractions.
    """
    findings: list = []
    try:
        space = chart._chartSpace
    except (AttributeError, ValueError):
        return findings
    C = "{%s}" % _CHART_NS
    for plot_tag in ("barChart", "lineChart"):
        for plot in space.iter(C + plot_tag):
            for ser in plot.findall(C + "ser"):
                # series display name
                nm = ser.find(".//" + C + "tx//" + C + "v")
                name = nm.text if nm is not None and nm.text else "?"
                # every formatCode attached to this series (data label + cache)
                codes = []
                dl = ser.find(C + "dLbls/" + C + "numFmt")
                if dl is not None and dl.get("formatCode"):
                    codes.append(dl.get("formatCode"))
                fc = ser.find(".//" + C + "val//" + C + "numCache/" + C + "formatCode")
                if fc is not None and fc.text:
                    codes.append(fc.text)
                if not any(_LITERAL_PCT_RE.search(c) for c in codes):
                    continue
                # cached values
                vals = []
                for pt in ser.findall(".//" + C + "val//" + C + "pt/" + C + "v"):
                    try:
                        vals.append(abs(float(pt.text)))
                    except (TypeError, ValueError):
                        continue
                nz = [v for v in vals if v != 0.0]
                if nz and max(nz) < 1.0:
                    findings.append(Finding(
                        "block", "chart-pct-format", loc,
                        f'percent series "{name}" uses a literal "%" format on '
                        f'sub-1.0 values (max {max(nz):.4g}) — if these are '
                        f'fractions (e.g. 0.0167) they render ~100x too small; '
                        f'use a true percent operator (0.0%) or confirm the '
                        f'values are already percentage-points'))
    return findings


def _overlap_ratio(s1, s2) -> float:
    try:
        boxes = [(s.left, s.top, s.width, s.height) for s in (s1, s2)]
    except (AttributeError, TypeError):
        return 0.0
    for box in boxes:
        if any(v is None for v in box):
            return 0.0
    (l1, t1, w1, h1), (l2, t2, w2, h2) = boxes
    ix = max(0, min(l1 + w1, l2 + w2) - max(l1, l2))
    iy = max(0, min(t1 + h1, t2 + h2) - max(t1, t2))
    inter = ix * iy
    if inter == 0:
        return 0.0
    smaller = min(w1 * h1, w2 * h2)
    return inter / smaller if smaller else 0.0


# ----------------------------------------------------------------------------
# Text-output linting (.md / .csv)
# ----------------------------------------------------------------------------
def lint_text_output(path: str) -> list[Finding]:
    findings: list[Finding] = []
    name = os.path.basename(path)

    if not re.search(r"\d{4}-\d{2}-\d{2}", name):
        findings.append(Finding("block", "filename", name,
                                "output filename must contain a YYYY-MM-DD date "
                                "(the analysis-period date, not today)"))

    text = _read_text(path)
    if text is None:
        return findings + [Finding("info", "read", path, "could not decode file")]

    for hexcode, why in FORBIDDEN_HEX.items():
        if re.search(hexcode, text, re.I):
            findings.append(Finding("block", "forbidden-color", name,
                                    f"#{hexcode} present — {why}"))
    for pat, why in FORBIDDEN_WORD_PATTERNS:
        if pat.search(text):
            findings.append(Finding("block", "forbidden-word", name,
                                    f'"{pat.search(text).group(0)}" — {why}'))

    # Attribution header (reports only — .md). CLAUDE.md: header metadata must
    # carry account, date range, data source, agent name/spec reference.
    if name.endswith(".md"):
        head = "\n".join(text.splitlines()[:30]).lower()
        if not (re.search(r"data source|source:", head) and re.search(r"agent", head)):
            findings.append(Finding("nit", "attribution", name,
                                    "report header missing a data-source + agent "
                                    "attribution line (see Output Format Conventions)"))

    # Currency-blend heuristic
    if re.search(r"\bCAD\b", text) and re.search(r"\bUSD\b", text):
        note = re.search(r"0\.73|conversion|converted|exchange rate|usd-equiv", text, re.I)
        if not note:
            findings.append(Finding("block", "currency", name,
                                    "both CAD and USD appear with no conversion-rate "
                                    "footnote — never blend currencies in one total"))
    return findings


def _read_text(path: str) -> str | None:
    for enc in ("utf-8-sig", "utf-16", "utf-8", "latin-1"):
        try:
            with open(path, encoding=enc) as fh:
                return fh.read()
        except (UnicodeError, UnicodeDecodeError):
            continue
    return None


# ----------------------------------------------------------------------------
# Raw-export reconciliation helper
# ----------------------------------------------------------------------------
def summarize_raw(path: str) -> list[Finding]:
    """Print filtered click/impression/spend totals from a raw export.

    Filters the Bing `Month == "Total"` footer row (and the Microsoft copyright
    row) before summing — naive sums double clicks/impressions (CLAUDE.md).
    Sniffs UTF-16LE (Google/Bing real exports) vs utf-8-sig (Meta/[DSP Platform]).
    """
    raw = _read_text(path)
    if raw is None:
        return [Finding("info", "read", path, "could not decode raw export")]

    # Locate the header row, then sniff the delimiter from that row. The real
    # header carries several KPI tokens; a report TITLE like "[Account: MedTech] Campaign
    # report" carries only one — so require >=2 distinct metric tokens. Scan a
    # generous window: Microsoft Ads exports have a 6-line preamble (+ a blank),
    # so the header can sit at line index 6 — `lines[:6]` missed it.
    lines = raw.splitlines()
    metric_tokens = ("clicks", "impr", "spend", "cost", "ctr", "conv", "cpc")
    start = 0
    for idx, line in enumerate(lines[:15]):
        low = line.lower()
        if sum(tok in low for tok in metric_tokens) >= 2:
            start = idx
            break
    header_line = lines[start] if lines else ""
    delim = "\t" if header_line.count("\t") > header_line.count(",") else ","
    reader = csv.DictReader(io.StringIO("\n".join(lines[start:])), delimiter=delim)

    def numcol(row, *names):
        for n in names:
            for k in row:
                if k and k.strip().lower() == n.lower():
                    v = (row[k] or "").replace(",", "").replace("$", "").replace("%", "").strip()
                    if v and v not in ("--", "-"):
                        try:
                            return float(v)
                        except ValueError:
                            return 0.0
        return 0.0

    clicks = impr = spend = 0.0
    rows = skipped = 0
    for row in reader:
        month = (row.get("Month") or "").strip()
        joined = " ".join(str(v) for v in row.values()).lower()
        if month.lower() == "total" or "all rights reserved" in joined:
            skipped += 1
            continue
        clicks += numcol(row, "Clicks")
        impr += numcol(row, "Impressions", "Impr.")
        spend += numcol(row, "Spend", "Cost", "Amount spent (USD)")
        rows += 1

    hint = ""
    if rows and clicks == 0 and impr == 0 and spend == 0:
        hint = (" — no metric columns found (a metric-free Search Terms export?); "
                "reconcile clicks/spend against the campaign / ad-group / keyword "
                "performance report instead")
    return [Finding("info", "reconcile", os.path.basename(path),
                    f"rows={rows} (skipped Total/footer={skipped}) | "
                    f"clicks={clicks:,.0f} impressions={impr:,.0f} "
                    f"spend={spend:,.2f} — compare against the report's quoted totals{hint}")]


# ----------------------------------------------------------------------------
# Keyword-dedup cross-check: never recommend adding a keyword the account
# already runs (the duplicate/no-op trap). [Account: MedSurg] 2026-06-17 — Codex caught
# net-new adds for `nasal decolonization` / `forced air warming` that were
# already live exact keywords.
# ----------------------------------------------------------------------------
def _normalize_kw(s: str) -> str:
    s = (s or "").strip().lower()
    s = s.strip('[]"’\'+ ')          # strip match-type wrappers: [exact] "phrase" +mod
    return re.sub(r"\s+", " ", s).strip()


_CH_ADD_RE = re.compile(r"keyword added\s*[\r\n]+\s*(\[[^\]]+\]|\"[^\"]+\"|[^\r\n]+)", re.I)


def _extract_keywords_from_change_history(raw_path: str) -> set:
    """Collect keywords ADDED per a Google/Bing change-history export's
    'Changes' column — catches keywords added before the performance
    export's window that still have zero impressions today, which the
    Search-keyword-column check below can never see (the [Account: MedSurg]
    2026-07-08 blind spot: a keyword added months ago with 0 impressions
    all quarter was wrongly recommended as a net-new add)."""
    raw = _read_text(raw_path)
    if raw is None:
        return set()
    lines = raw.splitlines()
    start = None
    for idx, line in enumerate(lines[:15]):
        low = line.lower()
        if "changes" in low and ("campaign" in low or "ad group" in low):
            start = idx
            break
    if start is None:
        return set()
    header_line = lines[start]
    delim = "\t" if header_line.count("\t") > header_line.count(",") else ","
    reader = csv.reader(io.StringIO("\n".join(lines[start:])), delimiter=delim)
    header = next(reader, [])
    change_idx = None
    for i, c in enumerate(header):
        if (c or "").strip().lower() == "changes":
            change_idx = i
            break
    if change_idx is None:
        return set()
    added = set()
    for row in reader:
        if len(row) <= change_idx:
            continue
        for m in _CH_ADD_RE.finditer(row[change_idx]):
            token = m.group(1).strip()
            if token.startswith("[") and "]" in token:
                token = token[1:token.index("]")]
            elif token.startswith('"') and token.count('"') >= 2:
                token = token[1:token.index('"', 1)]
            else:
                token = token.split(":", 1)[0]
            kw = _normalize_kw(token)
            if kw:
                added.add(kw)
    return added


def _extract_active_keywords(raw_path: str) -> set:
    """Collect the account's ACTIVE keywords from a raw export: the
    'Search keyword' (or 'Keyword') column of a campaigns/search-terms
    export, UNIONED with any keywords a change-history export shows as
    added (see _extract_keywords_from_change_history) — a keyword can be
    active with zero impressions in the current performance export."""
    active = _extract_keywords_from_change_history(raw_path)
    raw = _read_text(raw_path)
    if raw is None:
        return active
    lines = raw.splitlines()
    start = 0
    for idx, line in enumerate(lines[:15]):
        low = line.lower()
        if "search keyword" in low or re.search(r'(^|[\t,])"?keyword"?([\t,]|$)', low):
            start = idx
            break
    header_line = lines[start] if lines else ""
    delim = "\t" if header_line.count("\t") > header_line.count(",") else ","
    reader = csv.DictReader(io.StringIO("\n".join(lines[start:])), delimiter=delim)
    col = None
    for c in (reader.fieldnames or []):
        if c and c.strip().lower() in ("search keyword", "keyword"):
            col = c
            break
    if not col:
        return active
    for row in reader:
        kw = _normalize_kw(row.get(col, ""))
        if kw and kw not in ("--", "-"):
            active.add(kw)
    return active


def _csv_recommends_keywords(path: str) -> bool:
    text = _read_text(path)
    if not text:
        return False
    first = (text.splitlines() or [""])[0].lower()
    return "keyword" in first and ("action_type" in first or "action" in first)


def lint_keyword_dedup(csv_path: str, active: set) -> list[Finding]:
    """Flag keyword-ADD rows whose keyword already exists in the account."""
    findings: list[Finding] = []
    text = _read_text(csv_path)
    if text is None:
        return findings
    reader = csv.DictReader(io.StringIO(text))
    cols = {(c or "").strip().lower(): c for c in (reader.fieldnames or [])}
    kwcol = cols.get("keyword")
    actcol = cols.get("action_type") or cols.get("action")
    if not kwcol:
        return findings
    checked = dupes = 0
    for n, row in enumerate(reader, start=2):
        act = (row.get(actcol, "") if actcol else "").strip().upper()
        kw = _normalize_kw(row.get(kwcol, ""))
        if not kw or "ADD" not in act or "KEYWORD" not in act:
            continue
        checked += 1
        if kw in active:
            dupes += 1
            findings.append(Finding(
                "block", "keyword-dedup", f"{os.path.basename(csv_path)}:row {n}",
                f'recommends adding "{kw}" but it is ALREADY an active keyword in the '
                f'account — duplicate/no-op. Validate every add against the live keyword '
                f'set; carve to exact only when the term is currently served via broad/phrase.'))
    if checked and not dupes:
        findings.append(Finding(
            "info", "keyword-dedup", os.path.basename(csv_path),
            f"{checked} keyword-add row(s) checked vs {len(active)} active keywords — "
            f"none duplicate an existing keyword ✓"))
    return findings


# ----------------------------------------------------------------------------
def lint_path(path: str) -> list[Finding]:
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pptx":
        return lint_deck(path)
    if ext in (".md", ".csv", ".txt"):
        return lint_text_output(path)
    return [Finding("info", "skip", path, f"no linter for {ext} files")]


def main(argv=None) -> int:
    ap = argparse.ArgumentParser(description="Advisory QA linter (never blocks).")
    ap.add_argument("paths", nargs="*", help="deliverable files to lint")
    ap.add_argument("--raw", action="append", default=[],
                    help="raw export(s) to summarize for total reconciliation")
    args = ap.parse_args(argv)

    all_findings: list[tuple[str, list[Finding]]] = []
    for p in args.paths:
        if os.path.exists(p):
            all_findings.append((p, lint_path(p)))
        else:
            all_findings.append((p, [Finding("info", "missing", p, "file not found")]))
    for p in args.raw:
        if os.path.exists(p):
            all_findings.append((p, summarize_raw(p)))
        else:
            all_findings.append((p, [Finding("info", "missing", p, "file not found")]))

    # Cross-check keyword-ADD recommendations against the account's active keyword
    # set (requires --raw with a 'Search keyword'/'Keyword' column).
    if args.raw:
        active: set = set()
        for p in args.raw:
            if os.path.exists(p):
                active |= _extract_active_keywords(p)
        if active:
            for path, findings in all_findings:
                if path.lower().endswith(".csv") and os.path.exists(path) and _csv_recommends_keywords(path):
                    findings.extend(lint_keyword_dedup(path, active))

    total_block = 0
    for path, findings in all_findings:
        print(f"\n=== {path} ===")
        if not findings:
            print("  ✓ no findings")
            continue
        for f in sorted(findings, key=lambda x: {"block": 0, "nit": 1, "info": 2}[x.severity]):
            print(f.render())
            if f.severity == "block":
                total_block += 1

    print(f"\nSummary: {total_block} blocking-severity finding(s) "
          f"across {len(all_findings)} file(s). (Advisory — nothing was blocked.)")
    return 0  # always advisory


if __name__ == "__main__":
    sys.exit(main())

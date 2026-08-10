"""Canonical PPT building helpers for [Client]-branded executive decks.

Importable module backing the helper signatures documented in CLAUDE.md
("PPT Design Conventions"). Use these helpers from build scripts (e.g.,
the per-slide append scripts under `/tmp/append_slide_<n>.py` used to
assemble the HIS Q1 2026 deck) so the styling stays consistent across
runs and across analysts.

Usage:
    from scripts.ppt_helpers import (
        TEAL_DARK, TEAL_MID, TEAL_BRIGHT, CREAM_BG, GRAY_BODY, WHITE,
        set_title, add_subtitle_with_rule, add_section_header,
        add_textbox, add_bullets, render_table,
        add_card_light, add_card_dark, drop_slides_after,
    )

Design rules and color usage are documented in CLAUDE.md
"PPT Design Conventions (refined Apr 2026)". The reference implementation
this module preserves was used to build
`outputs/reports/ppt/executive_summary_his_blended_2026-04-26.pptx`.
"""

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# --- [Client] theme palette (sourced from templates/report_template.pptx theme1.xml) ---
TEAL_DARK    = RGBColor(0x01, 0x33, 0x2B)  # dk1   — primary text/title/headers
TEAL_MID     = RGBColor(0x19, 0xA5, 0x91)  # accent3 — secondary accents, card borders
TEAL_BRIGHT  = RGBColor(0x1C, 0xCF, 0x93)  # dk2   — divider rules only (never solid fill)
ACCENT_GREEN = RGBColor(0x05, 0xDD, 0x4D)  # accent1 — divider rules only
WARM_TAN     = RGBColor(0xC7, 0xA7, 0x9C)  # accent5 — fourth chart series only
CREAM_BG     = RGBColor(0xF7, 0xF4, 0xEA)  # soft callout body, light card backgrounds
GRAY_BODY    = RGBColor(0x5B, 0x5B, 0x5B)  # italic subtitles, footnotes
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)


# --- [Acquiring Co] theme palette (sourced from templates/acquiring_co_template.pptx theme1.xml) ---
# Parallel palette for the Purification & Filtration accounts that have moved
# [Client] -> [Acquiring Co] ([Account: MOEM], [Account: IWFB], BioPharma). Build on the TF template's
# own layouts and pass these colors into the helpers below (set_title color=,
# render_table header_fill=, add_card_* fill_color=, etc.). Never recolor a
# [Client]-teal deck into TF — the two palettes do not mix.
TF_NEAR_BLACK = RGBColor(0x22, 0x22, 0x22)  # dk1   — body text / titles
TF_RED        = RGBColor(0xE7, 0x13, 0x15)  # dk2   — primary brand accent (rules, emphasis)
TF_INDIGO     = RGBColor(0x26, 0x21, 0x60)  # accent1 — headers, table header bar, dark cards
TF_CYAN       = RGBColor(0x9A, 0xD3, 0xDC)  # accent2 — light accent, eyebrow on dark cards
TF_GOLD       = RGBColor(0xF0, 0xB2, 0x34)  # accent3 — chart series
TF_ORANGE     = RGBColor(0xEA, 0x76, 0x00)  # accent4 — chart series
TF_GRAY       = RGBColor(0x54, 0x58, 0x59)  # accent5 — italic subtitles, footnotes
TF_OLIVE      = RGBColor(0xB4, 0xBD, 0x01)  # accent6 — chart series
TF_CARD_BG    = RGBColor(0xEE, 0xF4, 0xF6)  # pale cyan tint — light card background


def set_title(slide, text, *, color=None, size=None, bold=None, font_name=None):
    """Fill the layout's native title placeholder with `text` (left-aligned).

    By default the placeholder inherits the master's title styling (the
    [Client] template renders dark teal). Pass `color`/`size`/`bold`/
    `font_name` to override explicitly — needed when building on a different
    brand master (e.g. [Acquiring Co]) so the title color is deterministic
    rather than inherited."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = text
            for p in ph.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                if color is None and size is None and bold is None and font_name is None:
                    continue
                for run in p.runs:
                    if color is not None: run.font.color.rgb = color
                    if size is not None: run.font.size = Pt(size)
                    if bold is not None: run.font.bold = bold
                    if font_name is not None: run.font.name = font_name
            return ph
    return None


def add_subtitle_with_rule(slide, left, top, width, text, *,
                           subtitle_size=12, rule_color=TEAL_BRIGHT,
                           rule_thickness=Pt(1.25)):
    """Italic gray subtitle with a thin teal divider rule directly beneath."""
    box = slide.shapes.add_textbox(left, top, width, Inches(0.34))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.04); tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0); tf.margin_bottom = Inches(0)
    tf.text = text
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    for run in p.runs:
        run.font.size = Pt(subtitle_size); run.font.italic = True
        run.font.color.rgb = GRAY_BODY; run.font.name = "Arial"

    rule_top = top + Inches(0.38)
    rule = slide.shapes.add_connector(1, left, rule_top, left + width, rule_top)
    rule.line.color.rgb = rule_color
    rule.line.width = rule_thickness
    return box, rule


def add_textbox(slide, left, top, width, height, text, *, size=12, bold=False,
                italic=False, color=TEAL_DARK, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP, font_name="Arial"):
    """Plain styled text box."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.04); tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    tf.text = text
    p = tf.paragraphs[0]; p.alignment = align
    for run in p.runs:
        run.font.size = Pt(size); run.font.bold = bold; run.font.italic = italic
        run.font.color.rgb = color; run.font.name = font_name
    return box


def add_section_header(slide, left, top, width, text, *, size=14, color=TEAL_DARK):
    """Bold section label (dark teal by default; pass `color` for another brand)."""
    return add_textbox(slide, left, top, width, Inches(0.32),
                       text, size=size, bold=True, color=color)


def add_bullets(slide, left, top, width, height, items, *, size=14,
                color=TEAL_DARK, line_spacing=1.30, marker="▪ ",
                marker_color=None):
    """Bullet list. `items` may contain plain strings or (label, body) tuples
    where `label` renders bold and `body` renders regular."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.04); tf.margin_right = Inches(0.04)
    mcol = marker_color or color
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing; p.alignment = PP_ALIGN.LEFT
        if isinstance(item, tuple):
            label, body = item
            r1 = p.add_run(); r1.text = marker
            r1.font.size = Pt(size); r1.font.color.rgb = mcol; r1.font.name = "Arial"
            r2 = p.add_run(); r2.text = label
            r2.font.size = Pt(size); r2.font.bold = True
            r2.font.color.rgb = color; r2.font.name = "Arial"
            r3 = p.add_run(); r3.text = " " + body
            r3.font.size = Pt(size); r3.font.color.rgb = color; r3.font.name = "Arial"
        else:
            run = p.add_run(); run.text = marker + item
            run.font.size = Pt(size); run.font.color.rgb = color; run.font.name = "Arial"
    return box


def style_table_cell(cell, *, fill=None, text_color=TEAL_DARK, size=10,
                     bold=False, align=PP_ALIGN.LEFT):
    """Apply [Client] table styling to a single cell."""
    if fill is not None:
        cell.fill.solid(); cell.fill.fore_color.rgb = fill
    tf = cell.text_frame
    tf.margin_left = Inches(0.10); tf.margin_right = Inches(0.10)
    tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
    tf.word_wrap = True
    for p in tf.paragraphs:
        p.alignment = align
        if not p.runs and p.text:
            run = p.add_run(); run.text = p.text
        for run in p.runs:
            run.font.size = Pt(size); run.font.bold = bold
            run.font.color.rgb = text_color; run.font.name = "Arial"


def render_table(slide, left, top, width, height, headers, rows, *,
                 header_fill=TEAL_DARK, header_text=WHITE,
                 body_text=TEAL_DARK, body_fill=WHITE,
                 first_col_bold=False,
                 header_size=10, body_size=10, col_aligns=None):
    """Clean [Client] table — dark-teal header bar over a clean white body
    (no alt-row stripes). Pass `first_col_bold=True` for label-style tables
    (e.g., metric in column 1, values across)."""
    n_cols = len(headers); n_rows = len(rows) + 1
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    tbl = table_shape.table
    aligns = col_aligns or [PP_ALIGN.LEFT] * n_cols
    for c, header in enumerate(headers):
        cell = tbl.cell(0, c); cell.text = str(header)
        style_table_cell(cell, fill=header_fill, text_color=header_text,
                         size=header_size, bold=True, align=aligns[c])
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = tbl.cell(r + 1, c); cell.text = str(val)
            is_label = first_col_bold and c == 0
            style_table_cell(cell, fill=body_fill, text_color=body_text,
                             size=body_size, bold=is_label, align=aligns[c])
    return tbl


def add_card_light(slide, left, top, width, height, *,
                   eyebrow=None, title=None, body=None,
                   eyebrow_size=10, title_size=18, body_size=12,
                   corner_pct=0.04,
                   border_color=TEAL_MID, fill_color=CREAM_BG,
                   eyebrow_color=TEAL_MID, title_color=TEAL_DARK,
                   body_color=TEAL_DARK):
    """Cream rounded-corner card with a thin teal border (light variant).
    Use for neutral callouts. Pair with `add_card_dark` to highlight the
    recommended option in an A/B comparison. Pass the `*_color`/`fill_color`
    overrides to render the card in another brand palette (e.g. [Acquiring Co])."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  left, top, width, height)
    card.adjustments[0] = corner_pct
    card.line.color.rgb = border_color
    card.line.width = Pt(0.75)
    card.fill.solid(); card.fill.fore_color.rgb = fill_color
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.30); tf.margin_right = Inches(0.30)
    tf.margin_top = Inches(0.22); tf.margin_bottom = Inches(0.22)
    tf.text = ""
    first = True
    if eyebrow:
        p = tf.paragraphs[0]; first = False
        r = p.add_run(); r.text = eyebrow.upper()
        r.font.size = Pt(eyebrow_size); r.font.bold = True
        r.font.color.rgb = eyebrow_color; r.font.name = "Arial"
    if title:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(2)
        r = p.add_run(); r.text = title
        r.font.size = Pt(title_size); r.font.bold = True
        r.font.color.rgb = title_color; r.font.name = "Arial"
    if body:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.space_before = Pt(6); p.line_spacing = 1.25
        r = p.add_run(); r.text = body
        r.font.size = Pt(body_size); r.font.color.rgb = body_color; r.font.name = "Arial"
    return card


def add_card_dark(slide, left, top, width, height, *,
                  eyebrow=None, title=None, body=None,
                  eyebrow_size=10, title_size=18, body_size=12,
                  corner_pct=0.04,
                  fill_color=TEAL_DARK, eyebrow_color=TEAL_BRIGHT,
                  title_color=WHITE, body_color=WHITE):
    """Solid dark rounded card (inverse variant), dark-teal with white text by
    default. Use to emphasize the recommended option in a paired comparison.
    Pass `fill_color`/`eyebrow_color`/... for another brand palette (e.g. a
    [Acquiring Co] indigo or red action card)."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  left, top, width, height)
    card.adjustments[0] = corner_pct
    card.line.fill.background()
    card.fill.solid(); card.fill.fore_color.rgb = fill_color
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.30); tf.margin_right = Inches(0.30)
    tf.margin_top = Inches(0.22); tf.margin_bottom = Inches(0.22)
    tf.text = ""
    first = True
    if eyebrow:
        p = tf.paragraphs[0]; first = False
        r = p.add_run(); r.text = eyebrow.upper()
        r.font.size = Pt(eyebrow_size); r.font.bold = True
        r.font.color.rgb = eyebrow_color; r.font.name = "Arial"
    if title:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(2)
        r = p.add_run(); r.text = title
        r.font.size = Pt(title_size); r.font.bold = True
        r.font.color.rgb = title_color; r.font.name = "Arial"
    if body:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.space_before = Pt(6); p.line_spacing = 1.25
        r = p.add_run(); r.text = body
        r.font.size = Pt(body_size); r.font.color.rgb = body_color; r.font.name = "Arial"
    return card


def drop_slides_after(prs, keep_n):
    """Remove all slides past index `keep_n` (0-based). Lets per-slide build
    scripts be re-run without producing duplicate slides — call this at the
    top of an append script with `keep_n=N-1` before `prs.slides.add_slide(...)`."""
    sldIdLst = prs.slides._sldIdLst
    slides = list(sldIdLst)
    for sld in slides[keep_n:]:
        rId = sld.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        prs.part.drop_rel(rId)
        sldIdLst.remove(sld)
